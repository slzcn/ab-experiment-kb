#!/usr/bin/env python3
# 文档 → Markdown 转码器（带图）。
# 把 PPT/Word/PDF 等常见文档抽成 markdown 文章：正文文字 + 内嵌图片。
# 图片模型沿用本项目最新那篇文章的做法：
#   - 抽出的图片写到 assets/<slug>-<n>.<ext>
#   - md 里用 raw.githubusercontent.com 的 URL 引用（线上/妙搭/本地都能显示）
#
# 用法：
#   python3 doc_to_md.py <文件路径> [--slug 前缀] [--title 标题] [--json]
#   默认输出 markdown 到 stdout；--json 输出 {title, md, images:[存好的文件名]}
#
# 依赖（本机已装）：python-pptx / python-docx / PyMuPDF(fitz) / Pillow
# 被 admin_server.py 的批量上传接口调用，也可命令行单独跑。

import sys, os, re, json, hashlib, io

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
sys.path.insert(0, HERE)
from kb_common import slugify as _slugify

# 图片线上引用前缀（与最新文章一致：raw.githubusercontent 指向仓库 assets/）
RAW_BASE = "https://raw.githubusercontent.com/slzcn/ab-experiment-kb/main/assets/"

MIN_IMG_BYTES = 3000      # 小于此的多半是图标/装饰，跳过
MIN_IMG_WH    = 80        # 任一边小于此像素跳过（去掉小logo/项目符号）


def _img_ok(data):
    """过滤太小或纯图标的图片，返回 (ok, ext)。"""
    if not data or len(data) < MIN_IMG_BYTES:
        return False, None
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(data))
        w, h = im.size
        if w < MIN_IMG_WH or h < MIN_IMG_WH:
            return False, None
        fmt = (im.format or "PNG").lower()
        ext = {"jpeg": "jpg"}.get(fmt, fmt)
        if ext not in ("jpg", "png", "gif", "webp"):
            ext = "png"
        return True, ext
    except Exception:
        return False, None


def _save_img(data, slug, seen):
    """存图到 assets/，去重（按内容 hash），返回线上 md 引用 URL 或 None。"""
    ok, ext = _img_ok(data)
    if not ok:
        return None
    h = hashlib.md5(data).hexdigest()[:8]
    if h in seen:
        return seen[h]                       # 同一张图复用已存文件
    os.makedirs(ASSETS, exist_ok=True)
    name = f"{slug}-{h}.{ext}"
    with open(os.path.join(ASSETS, name), "wb") as f:
        f.write(data)
    url = RAW_BASE + name
    seen[h] = url
    return url


# ---------------- PPTX：按幻灯片顺序，标题/正文层级 + 表格 + 该页图片 ----------------
def _shape_text_md(shp, is_title):
    """把一个文本框转成 md：标题框 → 小标题；正文框 → 保留项目符号层级的列表/段落。"""
    tf = shp.text_frame
    lines = []
    for p in tf.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        lvl = getattr(p, "level", 0) or 0
        lines.append((lvl, t))
    if not lines:
        return ""
    if is_title:
        return "## " + " ".join(t for _, t in lines)
    # 正文：多行→带缩进的无序列表（保留 PPT 的层级）；单行→普通段落
    if len(lines) == 1:
        return lines[0][1]
    return "\n".join(("  " * lvl) + "- " + t for lvl, t in lines)


def _pptx(path, slug):
    from pptx import Presentation
    prs = Presentation(path)
    seen = {}
    imgs = []
    md_parts = []
    BIG = 1 << 60  # top 缺失时排到末尾
    for slide in prs.slides:
        try:
            title_shp = slide.shapes.title
        except Exception:
            title_shp = None
        # 收集本页所有内容为 (top, kind, md)，再按纵向位置统一排序还原阅读顺序
        items = []
        for shp in slide.shapes:
            top = shp.top if getattr(shp, "top", None) is not None else BIG
            # 图片
            if shp.shape_type == 13:
                try:
                    url = _save_img(shp.image.blob, slug, seen)
                    if url: items.append((top, "img", url))
                except Exception:
                    pass
                continue
            # 表格
            if getattr(shp, "has_table", False) and shp.has_table:
                rows = [[c.text for c in r.cells] for r in shp.table.rows]
                tbl = _rows_to_md_table([r for r in rows if any((c or "").strip() for c in r)])
                if tbl: items.append((top, "tbl", tbl))
                continue
            # 文本框
            if shp.has_text_frame and shp.text_frame.text.strip():
                is_title = shp is title_shp
                md = _shape_text_md(shp, is_title)
                if md:
                    # 标题占位符永远置顶（top 记为 -1）
                    items.append((-1 if is_title else top, "title" if is_title else "body", md))
        items.sort(key=lambda x: x[0])
        # 若整页没有标题占位符，把最靠上的正文首行提升为小标题
        if items and not any(k == "title" for _, k, _ in items):
            for idx, (top, kind, md) in enumerate(items):
                if kind == "body":
                    first, *rest = md.split("\n")
                    first = first.lstrip("- ").strip()
                    new = [(top, "title", "## " + first)]
                    if rest and "\n".join(rest).strip():
                        new.append((top, "body", "\n".join(rest)))
                    items[idx:idx+1] = new
                    break
        parts = []
        for top, kind, md in items:
            if kind == "img":
                imgs.append(md); parts.append(f"![]({md})")
            else:
                parts.append(md)
        if parts:
            md_parts.append("\n\n".join(parts))
    return "\n\n".join(md_parts), imgs


# ---------------- DOCX：按文档流顺序穿插段落与图片 ----------------
def _docx(path, slug):
    import docx
    from docx.document import Document as _Doc
    d = docx.Document(path)
    seen = {}
    imgs = []
    # rId -> 图片二进制
    rels = {}
    for rid, rel in d.part.rels.items():
        if "image" in rel.reltype:
            try:
                rels[rid] = rel.target_part.blob
            except Exception:
                pass
    ns = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    md_parts = []
    for p in d.paragraphs:
        # 段落文字
        txt = p.text.strip()
        if txt:
            style = (p.style.name or "").lower()
            if "heading 1" in style or style == "title":
                md_parts.append(f"## {txt}")
            elif "heading" in style:
                md_parts.append(f"### {txt}")
            else:
                md_parts.append(txt)
        # 段落里的内嵌图
        blips = p._p.findall(".//a:blip", ns)
        for b in blips:
            rid = b.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
            if rid and rid in rels:
                url = _save_img(rels[rid], slug, seen)
                if url:
                    md_parts.append(f"![]({url})")
                    imgs.append(url)
    # 表格（转义 | 与换行，统一列数，避免撑破/错列）
    for tb in d.tables:
        rows = [[c.text for c in r.cells] for r in tb.rows if any(c.text.strip() for c in r.cells)]
        tbl = _rows_to_md_table(rows)
        if tbl:
            md_parts.append(tbl)
    return "\n\n".join(md_parts), imgs


# ---------------- PDF：按页抽文字 + 内嵌图 ----------------
def _pdf(path, slug):
    import fitz
    doc = fitz.open(path)
    seen = {}
    imgs = []
    md_parts = []
    for pno in range(len(doc)):
        page = doc[pno]
        txt = page.get_text().strip()
        if txt:
            md_parts.append(txt)
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n - pix.alpha >= 4:          # CMYK → RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                data = pix.tobytes("png")
                url = _save_img(data, slug, seen)
                if url:
                    md_parts.append(f"![]({url})")
                    imgs.append(url)
            except Exception:
                pass
    return "\n\n".join(md_parts), imgs


# ---------------- 纯文本 / 表格：自包含处理（无图，不依赖外部脚本，Action 里可跑）----------------
def _plain(path):
    return open(path, encoding="utf-8", errors="replace").read().strip()


def _csv_md(path):
    import csv
    rows = []
    with open(path, encoding="utf-8", errors="replace", newline="") as f:
        for r in csv.reader(f):
            if any((c or "").strip() for c in r):
                rows.append(r)
    return _rows_to_md_table(rows)


def _fmt_cell(v):
    """把单元格值格式化成适合 md 表格的一行文本：保留位置、转义 | 、把换行变 <br>。"""
    if v is None:
        return ""
    if isinstance(v, float):
        # 整数值的浮点去掉 .0（Excel 里 9 常被读成 9.0）
        v = int(v) if v.is_integer() else round(v, 6)
    s = str(v).strip()
    s = s.replace("\\", "\\\\").replace("|", "\\|")   # 转义竖线，避免撑破列
    s = re.sub(r"\r?\n", "<br>", s)                    # 单元格内换行 → <br>，保证一行
    return s


def _rows_to_md_table(rows):
    """rows: 二维列表(已按最大列数对齐)。生成合法 md 表格：首行表头 + 分隔行 + 数据行。"""
    if not rows:
        return ""
    ncol = max(len(r) for r in rows)
    def line(cells):
        cells = list(cells) + [""] * (ncol - len(cells))   # 补齐到统一列数
        return "| " + " | ".join(_fmt_cell(c) for c in cells) + " |"
    out = [line(rows[0]), "|" + "---|" * ncol]
    out += [line(r) for r in rows[1:]]
    return "\n".join(out)


def _xlsx_md(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        # 收集非空行（保留单元格位置，不丢中间的空格），并裁掉整行/整列皆空的边缘
        raw = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None and str(c).strip() for c in row):
                raw.append(list(row))
        if not raw:
            continue
        # 裁掉尾部全空的列（避免一堆空列）
        maxc = 0
        for r in raw:
            for j in range(len(r) - 1, -1, -1):
                if r[j] is not None and str(r[j]).strip():
                    maxc = max(maxc, j + 1); break
        rows = [r[:maxc] for r in raw]
        tbl = _rows_to_md_table(rows)
        if tbl:
            title = ws.title if ws.title and not ws.title.lower().startswith("sheet") else ""
            out.append((f"## {title}\n\n" if title else "") + tbl)
    return "\n\n".join(out)


def _html_text(path):
    raw = open(path, encoding="utf-8", errors="replace").read()
    raw = re.sub(r"<(script|style)[^>]*>[\s\S]*?</\1>", "", raw, flags=re.I)
    txt = re.sub(r"<[^>]+>", " ", raw)
    txt = re.sub(r"&nbsp;", " ", txt)
    return re.sub(r"[ \t]*\n[ \t]*", "\n", re.sub(r"[ \t]+", " ", txt)).strip()


def convert(path, slug=None, title=None):
    ext = os.path.splitext(path)[1].lower()
    base = os.path.splitext(os.path.basename(path))[0]
    title = title or base
    slug = _slugify(slug or base)
    if ext == ".pptx":
        md, imgs = _pptx(path, slug)
    elif ext == ".docx":
        md, imgs = _docx(path, slug)
    elif ext == ".pdf":
        md, imgs = _pdf(path, slug)
    elif ext in (".txt", ".md", ".markdown"):
        md, imgs = _plain(path), []
    elif ext == ".csv":
        md, imgs = _csv_md(path), []
    elif ext in (".xlsx", ".xlsm"):
        md, imgs = _xlsx_md(path), []
    elif ext in (".html", ".htm"):
        md, imgs = _html_text(path), []
    else:
        # 老二进制格式（.doc .ppt .xls .rtf）需 LibreOffice/textutil，Action 环境不具备
        raise RuntimeError(f"暂不支持的格式 {ext}（请先另存为 .docx/.pptx/.xlsx/.pdf 再上传）")
    md = re.sub(r"\n{3,}", "\n\n", md).strip()
    return {"title": title, "md": md, "images": imgs}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python3 doc_to_md.py <文件> [--slug 前缀] [--title 标题] [--json]")
        sys.exit(1)
    path = sys.argv[1]
    as_json = "--json" in sys.argv

    def _opt(flag):
        if flag in sys.argv:
            i = sys.argv.index(flag)
            if i + 1 < len(sys.argv):
                return sys.argv[i + 1]
        return None

    if not os.path.exists(path):
        print("文件不存在:", path); sys.exit(1)
    try:
        r = convert(path, slug=_opt("--slug"), title=_opt("--title"))
    except Exception as e:
        print(f"转码失败({type(e).__name__}): {e}"); sys.exit(2)
    if as_json:
        print(json.dumps(r, ensure_ascii=False))
    else:
        print(f"# {r['title']}\n")
        print(r["md"])
        print(f"\n\n[抽出图片 {len(r['images'])} 张 → assets/]", file=sys.stderr)
